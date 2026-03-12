from pptx import Presentation
from langchain_core.documents import Document

def load_pptx(file_path):
    """
    Loads a .pptx file and extracts text from all slides.
    """
    try:
        prs = Presentation(file_path)
        documents = []

        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            content = "\n".join(slide_text)
            
            if content.strip():
                meta = {
                    "source": file_path,
                    "slide_number": i + 1
                }
                documents.append(Document(page_content=content, metadata=meta))

        return documents

    except Exception as e:
        print(f"Error loading PPTX {file_path}: {e}")
        return []
